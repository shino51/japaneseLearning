from pathlib import Path
import csv
import urllib.request

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent
CSV_FILE = ROOT / "data" / "pokemon.csv"
IMAGE_DIR = ROOT / "images"
OUTPUT = ROOT / "pokemon_zukan.pptx"

A4_W = Cm(21.0)
A4_H = Cm(29.7)

# PokeAPI sprites repository: official artwork PNGs.
ARTWORK_URL = (
    "https://raw.githubusercontent.com/PokeAPI/sprites/master/"
    "sprites/pokemon/other/official-artwork/{dex_id}.png"
)


def download_missing_images(pokemon):
    IMAGE_DIR.mkdir(exist_ok=True)
    for p in pokemon:
        target = IMAGE_DIR / p["image"]
        if target.exists():
            continue
        url = ARTWORK_URL.format(dex_id=int(p["id"]))
        print(f"Downloading {p['name_ja']} ...")
        urllib.request.urlretrieve(url, target)


def add_text(slide, text, left, top, width, height, size, bold=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "BIZ UDPGothic"
    r.font.size = Pt(size)
    r.font.bold = bold
    return box


def add_pokemon_slide(prs, pokemon_group):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # A4 portrait: 2 columns x 2 rows = 4 Pokemon per page.
    margin_x = Inches(0.45)
    margin_top = Inches(0.35)
    margin_bottom = Inches(0.35)
    gap_x = Inches(0.30)
    gap_y = Inches(0.30)
    card_w = (A4_W - margin_x * 2 - gap_x) / 2
    card_h = (A4_H - margin_top - margin_bottom - gap_y) / 2

    for i, p in enumerate(pokemon_group):
        col = i % 2
        row = i // 2
        left = margin_x + col * (card_w + gap_x)
        top = margin_top + row * (card_h + gap_y)

        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
        )
        sh.fill.background()
        sh.line.width = Pt(0.8)

        add_text(slide, f"#{int(p['id']):03d}",
                 left + Inches(0.08), top + Inches(0.37),
                 Inches(0.55), Inches(0.25), 9, True)
        # Image area: square, centered in the card
        image_path = IMAGE_DIR / p["image"]

        img_size = min(card_w - Inches(0.30), Inches(2.15))

        img_top = top + Inches(0.75)

        if image_path.exists():
            # Keep the original aspect ratio and center the image.
            from PIL import Image

            with Image.open(image_path) as im:
                image_w, image_h = im.size

            if image_w >= image_h:
                display_w = img_size
                display_h = img_size * image_h / image_w
            else:
                display_h = img_size
                display_w = img_size * image_w / image_h

            display_left = left + (card_w - display_w) / 2
            display_top = img_top + (img_size - display_h) / 2

            slide.shapes.add_picture(
                str(image_path),
                display_left,
                display_top,
                width=display_w,
                height=display_h
            )
        else:
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left + Inches(0.15),
                img_top,
                img_size,
                img_size
            )
            placeholder.fill.background()
            placeholder.text_frame.text = f"[画像なし]\n{p['image']}"

        add_text(slide, p["name_de"],
                 left + Inches(0.12), top + Inches(3.25), # 2.55
                 card_w - Inches(0.24), Inches(0.30), 14)
        add_text(slide, " ".join(list(p["hiragana"])),
                 left + Inches(0.10), top + Inches(3.70),
                 card_w - Inches(0.20), Inches(0.42), 24, True)
        add_text(slide, p["name_ja"],
                 left + Inches(0.10), top + Inches(4.30),
                 card_w - Inches(0.20), Inches(0.35), 15)

        types = " / ".join(x for x in (p["type1"], p["type2"]) if x)
        if types:
            add_text(slide, f"タイプ：{types}",
                     left + Inches(0.10), top + Inches(4.70),
                     card_w - Inches(0.20), Inches(0.25), 9)

        if p["encountered_date"]:
            add_text(slide, f"出会った日：{p['encountered_date']}",
                     left + Inches(0.10), top + Inches(4.02),
                     card_w - Inches(0.20), Inches(0.22), 7)


def main():
    with CSV_FILE.open("r", encoding="utf-8-sig", newline="") as f:
        pokemon = list(csv.DictReader(f))

    # Images are downloaded automatically the first time.
    download_missing_images(pokemon)

    prs = Presentation()
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    for i in range(0, len(pokemon), 4):
        add_pokemon_slide(prs, pokemon[i:i+4])

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT}")
    print(f"Pages: {len(pokemon)}")


if __name__ == "__main__":
    main()
