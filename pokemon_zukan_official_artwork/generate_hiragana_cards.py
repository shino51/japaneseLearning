from __future__ import annotations

import argparse
import csv
import random
from collections import Counter
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
CSV_PATH = ROOT / "data" / "pokemon.csv"
OUTPUT = ROOT / "hiragana_cards.pptx"

# Cards are laid out for A4 printing.
A4_W = Inches(8.2677)
A4_H = Inches(11.6929)

CARD_SIZE = Inches(1.7)
CARD_GAP_X = Inches(0)
CARD_GAP_Y = Inches(0)
MARGIN_X = Inches(0.55)
MARGIN_Y = Inches(0.55)

START_DATE = "2026-08-01"
END_DATE = "2026-08-16"

# ダミーのひらがなカードを何枚追加するか
DISTRACTOR_COUNT = 0


# Characters that are useful distractors for this kind of game.
DISTRACTOR_CHARS = list(
    "あいうえお"
    "かきくけこ"
    "がぎぐげご"
    "さしすせそ"
    "ざじずぜぞ"
    "たちつてと"
    "だぢづでど"
    "なにぬねの"
    "はひふへほ"
    "ばびぶべぼ"
    "ぱぴぷぺぽ"
    "まみむめも"
    "やゆよ"
    "らりるれろ"
    "わをん"
    "っー"
)


def remove_shadow(shape):
    """PowerPointの図形・文字に設定された影をXMLから削除する。"""
    for element in shape._element.iter():
        for effect_list in element.findall(".//" + qn("a:effectLst")):
            parent = effect_list.getparent()
            if parent is not None:
                parent.remove(effect_list)

def parse_date(value: str) -> date | None:
    if not value:
        return None
    return date.fromisoformat(value)


def load_pokemon(start: date, end: date) -> list[dict]:
    with CSV_PATH.open("r", encoding="utf-8-sig", newline="") as f:
        rows = list(csv.DictReader(f))

    selected = []
    for row in rows:
        d = parse_date(row.get("encountered_date", ""))
        if d is not None and start <= d <= end:
            selected.append(row)

    return selected


def add_card(slide, char: str, x, y):
    shape = slide.shapes.add_shape(
        1, x, y, CARD_SIZE, CARD_SIZE  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.background()
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = char
    r.font.name = "BIZ UDPGothic"
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0, 0, 0)

    # 図形そのものの影を削除
    remove_shadow(shape)


def generate_cards(chars: list[str], output: Path):
    prs = Presentation()
    prs.slide_width = A4_W
    prs.slide_height = A4_H

    per_row = 4
    per_col = 6
    per_page = per_row * per_col

    for page_start in range(0, len(chars), per_page):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        page_chars = chars[page_start:page_start + per_page]

        for i, char in enumerate(page_chars):
            row = i // per_row
            col = i % per_row
            x = MARGIN_X + col * (CARD_SIZE + CARD_GAP_X)
            y = MARGIN_Y + row * (CARD_SIZE + CARD_GAP_Y)
            add_card(slide, char, x, y)

    prs.save(output)


def main():
    start = date.fromisoformat(START_DATE)
    end = date.fromisoformat(END_DATE)

    if start > end:
        raise SystemExit("START_DATE は END_DATE 以下にしてください。")

    pokemon = load_pokemon(start, end)

    if not pokemon:
        raise SystemExit(
            f"{start} ～ {end} に出会ったポケモンがCSVにありません。"
            " encountered_date を確認してください。"
        )

    # 指定期間の全ポケモンについて必要な文字を集計
    required = {}

    for p in pokemon:
        char_counts = Counter(p["hiragana"])

        for char, count in char_counts.items():
            required[char] = max(required.get(char, 0), count)

    cards = []
    for char, count in required.items():
        cards.extend([char] * count)


    # ダミー文字を追加
    if DISTRACTOR_COUNT != 0:
        rng = random.Random()
        distractors = [
            rng.choice(DISTRACTOR_CHARS)
            for _ in range(DISTRACTOR_COUNT)
        ]
        cards += distractors
        print(f"ダミーカード: {len(distractors)}枚 追加")

    print("対象ポケモン:")
    for p in pokemon:
        print(f"  {p['name_ja']} ({p['hiragana']})")

    print(f"必要な文字カード: {sum(required.values())}枚")
    print(f"合計: {len(cards)}枚")

    generate_cards(cards, OUTPUT)
    print(f"出力: {OUTPUT}")


if __name__ == "__main__":
    main()
